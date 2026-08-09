#!/usr/bin/env python3
"""Generate PowerPoint: EfficientAI Call Import Architecture & Scaling Guide."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "presentations" / "EfficientAI_Call_Import_Architecture_and_Scaling.pptx"

# Brand-ish palette
NAVY = RGBColor(0x1A, 0x36, 0x5D)
TEAL = RGBColor(0x00, 0x96, 0x88)
SLATE = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT_ORANGE = RGBColor(0xEA, 0x58, 0x0C)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, NAVY)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.0), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    sub = slide.shapes.add_textbox(Inches(0.6), Inches(3.6), Inches(12.0), Inches(1.2))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(20)
    sp.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    sp.alignment = PP_ALIGN.LEFT


def _add_section_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, TEAL)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT


def _add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    *,
    subtitle: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.7))
    tfp = title_box.text_frame.paragraphs[0]
    tfp.text = title
    tfp.font.size = Pt(28)
    tfp.font.bold = True
    tfp.font.color.rgb = NAVY

    y = 1.15
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.3), Inches(0.5))
        sfp = sub_box.text_frame.paragraphs[0]
        sfp.text = subtitle
        sfp.font.size = Pt(14)
        sfp.font.italic = True
        sfp.font.color.rgb = SLATE
        y += 0.55

    body = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(12.2), Inches(6.5 - y))
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(16 if len(bullet) < 120 else 14)
        p.font.color.rgb = SLATE
        p.space_after = Pt(8)


def _add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.7))
    tfp = title_box.text_frame.paragraphs[0]
    tfp.text = title
    tfp.font.size = Pt(26)
    tfp.font.bold = True
    tfp.font.color.rgb = NAVY

    nrows = len(rows) + 1
    ncols = len(headers)
    table_shape = slide.shapes.add_table(nrows, ncols, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.35 * nrows + 0.3))
    table = table_shape.table

    col_width = Inches(12.5 / ncols)
    for c in range(ncols):
        table.columns[c].width = int(col_width)

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = SLATE
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title ---
    _add_title_slide(
        prs,
        "EfficientAI Call Import Architecture & Scaling",
        "Enterprise deployment guide  |  Sharded Postgres  |  Redis fair-share  |  Kubernetes + KEDA",
    )

    # --- Agenda ---
    _add_content_slide(
        prs,
        "Agenda",
        [
            "Platform architecture overview",
            "Call import end-to-end pipeline (import → transcribe → eval)",
            "Postgres catalog + data shards vs Redis coordination",
            "Fair dispatch, inflight limits, and batch sizing",
            "Customer scenario: dual 10k workspace runs",
            "Before vs recommended configuration numbers",
            "Database connection budget (without PgBouncer)",
            "Kubernetes autoscaling with Prometheus + KEDA",
            "Future scope: PgBouncer / RDS Proxy",
            "Monitoring, load-test exit criteria, and next steps",
        ],
    )

    # --- Architecture ---
    _add_section_slide(prs, "1. Platform Architecture")

    _add_content_slide(
        prs,
        "High-Level Architecture",
        [
            "EfficientAI runs as Kubernetes services: API, worker-imports, worker (audio-metrics), Redis, and Postgres.",
            "Catalog DB (efficientai_catalog): metadata — CallImport headers, CallImportEvaluation, org/workspace, metrics catalog, shard registry.",
            "Data shards (efficientai_data_01 … _04): high-volume row data — CallImportRow, CallImportEvaluationRow.",
            "Redis: Celery broker/result backend + fair-share inflight counters + dispatch cursors + rate-limit state.",
            "Blob storage (S3/GCS/Azure): CSV uploads and call recordings.",
            "External APIs: STT/LLM providers, telephony recording URLs (Exotel, Plivo, direct URLs).",
        ],
        subtitle="1 catalog + 4 data shards (customer topology)",
    )

    _add_content_slide(
        prs,
        "Worker Topology",
        [
            "worker-imports (thread pool): queues imports, diarization, eval-control, evaluations",
            "  → I/O bound: recording fetch, STT, LLM diarisation, LLM metric scoring, fair dispatch",
            "  → Typical K8s: 8–16 replicas × 16 concurrency (recommended)",
            "worker (prefork): queues celery + audio-metrics",
            "  → CPU/audio bound: Praat, UTMOS, torch-based qualitative voice metrics",
            "  → Separate pool avoids OMP/torch deadlocks with high Celery concurrency",
            "Queue drain order (worker-imports): imports → diarization → eval-control → evaluations",
        ],
    )

    # --- Pipeline ---
    _add_section_slide(prs, "2. Call Import Pipeline")

    _add_content_slide(
        prs,
        "End-to-End Pipeline Stages",
        [
            "1. Upload CSV → API creates CallImport header (catalog) + stores CSV in blob storage",
            "2. Materialize rows → bulk_insert_mappings_on_shards() → CallImportRow on data shards",
            "3. Bulk import (optional phase) → fair_import_dispatch → process_call_import_row (fetch recording → S3)",
            "4. Create evaluation → CallImportEvaluation header (catalog) + materialize CallImportEvaluationRow on shards",
            "5. Fair eval dispatch → per row: import (if needed) → transcribe → audio metrics → LLM eval",
            "6. Rollup → parent counters and evaluation status updated; insights tasks may follow",
        ],
        subtitle="Unified eval pipeline: one Redis eval slot can cover import + transcribe + eval for a row",
    )

    _add_content_slide(
        prs,
        "Celery Queues & Key Tasks",
        [
            "imports — process_call_import_row, dispatch_fair_import_rows, bulk materialize/delete",
            "diarization — transcribe_call_import_row, dispatch_fair_diarization_rows",
            "eval-control — materialize/cancel/retry evaluation (operator actions, not head-of-line blocked)",
            "evaluations — dispatch_fair_eval_rows, evaluate_call_import_row, insights generation",
            "audio-metrics — evaluate_call_import_row_audio (separate worker service)",
            "Task time limits: transcribe soft 12 min; LLM eval soft 8 min; global Celery hard limit 30 min",
        ],
    )

    _add_content_slide(
        prs,
        "Per-Row Dispatch Decision Tree",
        [
            "_try_dispatch_single_row() (eval_dispatch.py) picks the next stage for each pending eval row:",
            "  • No S3 recording yet → enqueue process_call_import_row on imports (uses eval slot in eval chain)",
            "  • Needs diarised transcript → enqueue transcribe_call_import_row on diarization",
            "  • Audio-only metrics configured → enqueue evaluate_call_import_row_audio on audio-metrics",
            "  • Otherwise → enqueue evaluate_call_import_row on evaluations (LLM/comparison metrics)",
            "Slot held from dispatch until row completes → finish_eval_work_and_redispatch() releases slot",
            "Bulk CSV import (pre-eval) uses separate import:* Redis slots via fair_import_dispatch",
        ],
    )

    # --- Postgres vs Redis ---
    _add_section_slide(prs, "3. Postgres vs Redis")

    _add_content_slide(
        prs,
        "Division of Responsibility",
        [
            "Postgres = source of truth for durable state (rows, statuses, scores, headers, registry)",
            "Redis = ephemeral coordination (inflight caps, fair-scheduling cursors, dedupe locks, progress hashes)",
            "Celery broker (Redis): task messages — but bulk pending work lives in Postgres, not the queue",
            "Key insight: queue depth alone is a poor autoscaling signal — fair dispatch only enqueues when a slot is free",
        ],
    )

    _add_content_slide(
        prs,
        "What Lives in Postgres",
        [
            "Catalog DB:",
            "  • CallImport, CallImportEvaluation (headers, config, status)",
            "  • call_import_shard_slices registry (rebalance overrides)",
            "  • Organizations, workspaces, metrics, AI provider credentials",
            "Data shards (per shard):",
            "  • CallImportRow — recording_url, recording_s3_key, transcripts, import status",
            "  • CallImportEvaluationRow — eval status, metric_scores, celery_task_id",
            "Pending work query: status=pending AND celery_task_id IS NULL (scatter-gather across shards)",
        ],
    )

    _add_content_slide(
        prs,
        "What Lives in Redis",
        [
            "Eval inflight counters: eval:inflight:global | :org:{id} | :workspace:{id} | :job:{eval_id}",
            "Import inflight counters: import:inflight:global | :org:{id} | :workspace:{id}",
            "Slot task maps: eval:slot:task:{celery_task_id}, import:slot:task:{celery_task_id}",
            "Fair dispatch cursors: eval:fair:rr_cursor, eval:fair:rr_cursor:ws:{workspace_id} (and import equivalents)",
            "Dispatch locks/dedupe: eval:fair:dispatch_lock, eval:fair:dispatch_dedupe (15s backoff at capacity)",
            "Retry metadata: eval:restricted:row:{id}, eval:transcribe_overwrite:{evaluation_id}",
            "Telephony rate limits: telephony:import:credits:{fingerprint} (default 1000/min per credential)",
        ],
    )

    _add_content_slide(
        prs,
        "Data Flow: Postgres ↔ Redis ↔ Workers",
        [
            "① API writes headers + rows to Postgres (catalog + shards via router)",
            "② API/worker schedules dispatch_fair_* Celery task (message → Redis broker)",
            "③ Dispatcher reads pending rows from Postgres (scatter-gather on shards + catalog joins)",
            "④ For each row: acquire_*_slot() atomically increments Redis inflight counters (Lua script)",
            "⑤ If slot acquired → enqueue row task to Celery; store celery_task_id on shard row (Postgres)",
            "⑥ Worker executes task: reads/writes Postgres (catalog + one shard), calls external APIs",
            "⑦ Task completes → release_*_slot() decrements Redis counters → schedule next fair dispatch turn",
            "At capacity: dispatch stops enqueueing; pending rows remain in Postgres (queue may look empty)",
        ],
        subtitle="Pending backlog is in Postgres; Redis tracks who is allowed to run right now",
    )

    # --- Sharding ---
    _add_section_slide(prs, "4. Database Sharding")

    _add_content_slide(
        prs,
        "Shard Routing",
        [
            "Enabled via database.sharding.enabled: true in config.yml",
            "slice_id = row_index // row_chunk_size  (default row_chunk_size = 500)",
            "shard_id = SHA256(call_import_id : slice_id) mod N  (N = number of shards)",
            "Registry table call_import_shard_slices can override routing during rebalance",
            "10k row import → ~20 slices (500 rows each) → hash-distributed across 4 shards (~2.5k rows/shard)",
            "Router: app/db_sharding/router.py (ShardRouter)  |  Pools: app/db_sharding/pool_manager.py",
        ],
    )

    _add_content_slide(
        prs,
        "Connection Pools (Per Process)",
        [
            "Each worker/API process holds SQLAlchemy pools for catalog + every shard",
            "With 4 shards, pool_manager reduces per-shard pool but enforces floor: 8+8 = 16 connections/shard/process",
            "Catalog pool uses full pool_size + max_overflow (not divided)",
            "Each active task typically holds: 1 catalog session + 1 shard session (seconds to minutes)",
            "Fair dispatch may open ShardSessionCache: up to 4 shard sessions + 1 catalog per dispatch pass",
            "Rule: concurrency per pod ≤ catalog pool max — or threads block on pool timeout",
        ],
    )

    # --- Fair dispatch ---
    _add_section_slide(prs, "5. Fair Dispatch & Limits")

    _add_content_slide(
        prs,
        "Fair Dispatch Mechanics",
        [
            "Two-level round-robin:",
            "  • Global: rotate across workspaces that have pending rows (eval:fair:rr_cursor)",
            "  • Per-workspace: rotate across evaluations (eval:fair:rr_cursor:ws:{workspace_id})",
            "max_workspace_turns:",
            "  • 999 on eval create / catch-up → fill capacity across workspaces quickly",
            "  • 1 after each row completes → fair refill, one workspace turn at a time",
            "eval_fair_dispatch_batch_size: max rows attempted per workspace turn (default 75)",
            "  → Should align with eval_workspace_inflight_limit to avoid stair-step ramp",
        ],
    )

    _add_content_slide(
        prs,
        "Inflight Limit Hierarchy",
        [
            "Eval slot acquisition checks (all must pass):",
            "  workspace_limit → org_limit → global_limit → job_limit (per evaluation id)",
            "Import slot acquisition (bulk CSV only): workspace → org → global",
            "Effective parallel rows = min(job, workspace, org, global, worker_threads)",
            "Customer issue identified: eval_job_inflight_limit defaults to 75 — caps large single evaluations",
            "Even with 512 global and 28 concurrency, effective parallelism was ~75 rows for 5k/10k runs",
        ],
    )

    # --- Customer scenario ---
    _add_section_slide(prs, "6. Customer Scenario & Metrics")

    _add_content_slide(
        prs,
        "Observed Customer Baseline",
        [
            "Topology: 1 catalog Postgres + 4 data shard Postgres instances",
            "Deployment: Kubernetes with Prometheus + KEDA autoscaling",
            "Workload: 5,000 calls per workspace — full pipeline (import + transcribe + eval) in ~30 minutes",
            "Target: 10,000 calls in 30 minutes (single workspace) or dual 10k across two workspaces",
            "Original config: 8 pods × 28 concurrency, eval_global=512, eval_workspace=150, import_global=96",
            "Hidden bottleneck: eval_job_inflight_limit = 75 (default) — not explicitly raised",
        ],
    )

    _add_table_slide(
        prs,
        "Before vs Recommended Configuration",
        ["Setting", "Before (customer)", "Now (recommended)"],
        [
            ["KEDA min / max replicas", "8 / 16", "8 / 16"],
            ["worker_imports_concurrency", "28", "16"],
            ["Thread capacity @ max pods", "448", "256"],
            ["database pool_size / max_overflow", "10 / 15", "6 / 10"],
            ["Catalog pool max / pod", "25", "16 (= concurrency)"],
            ["eval_global_inflight_limit", "512", "256"],
            ["eval_workspace_inflight_limit", "150", "128"],
            ["eval_job_inflight_limit", "75 (default)", "128"],
            ["eval_fair_dispatch_batch_size", "75 (default)", "128"],
            ["import_global_inflight_limit", "96", "96 (keep)"],
            ["import_workspace_inflight_limit", "48", "48 (keep)"],
            ["import_fair_dispatch_batch_size", "75 (default)", "48"],
        ],
    )

    _add_table_slide(
        prs,
        "Effective Parallelism & Throughput (@ ~24 s/row)",
        ["Scenario", "Before (effective)", "Recommended @ max (16 pods)"],
        [
            ["Single 10k eval @ 8 pods", "~75 parallel (~30 min)", "~128 parallel (~31 min)"],
            ["Single 10k eval @ 16 pods", "~75 parallel (~30 min)", "~256 parallel (~16 min)"],
            ["Dual 10k (2 workspaces) @ 16 pods", "~150 parallel (2×75 job)", "~256 parallel (2×128 ws)"],
            ["Bulk import per workspace", "48 parallel fetches", "48 parallel (unchanged)"],
        ],
    )

    _add_content_slide(
        prs,
        "Why Recommended Numbers Change",
        [
            "Raise eval_job 75 → 128: removes hidden cap that prevented scaling past ~75 rows/eval",
            "Lower concurrency 28 → 16: catalog pool (25) < 28 threads caused pool timeout / DB throttling",
            "Align eval batch 75 → 128: one workspace turn can fill quota (no stair-step with ws limit 128)",
            "Keep import 96/48: bulk import phase stays fast for dual 10k; only fix import batch 75 → 48",
            "eval_global 512 → 256: match realistic max capacity (16 × 16) without over-dispatching DB",
            "Fair share for 2 workspaces: eval_workspace = eval_global ÷ 2 = 128",
        ],
    )

    _add_table_slide(
        prs,
        "Fair-Share Limit Formulas",
        ["Limit", "Formula", "Example (2 heavy workspaces, max 16×16)"],
        [
            ["eval_global", "max_replicas × concurrency", "16 × 16 = 256"],
            ["eval_workspace", "eval_global ÷ N workspaces", "256 ÷ 2 = 128"],
            ["eval_job", "≥ eval_workspace", "128"],
            ["eval_fair_dispatch_batch_size", "= eval_workspace", "128"],
            ["import_global", "keep or import_ws × N", "48 × 2 = 96"],
            ["import_workspace", "import_global ÷ N", "96 ÷ 2 = 48"],
            ["import_fair_dispatch_batch_size", "= import_workspace", "48"],
        ],
    )

    # --- DB connections ---
    _add_section_slide(prs, "7. Database Connection Budget")

    _add_content_slide(
        prs,
        "Connection Math (No PgBouncer)",
        [
            "Assume: 16 worker-imports + 3 API pods = 19 processes at max scale",
            "Catalog connections: 19 × (pool_size + max_overflow) = 19 × 16 = ~304",
            "Per shard connections: 19 × 16 (shard pool floor) = ~304 per shard RDS",
            "Minimum RDS max_connections: catalog ≥ 450, each shard ≥ 450 (with admin headroom)",
            "Load-test exit criteria (docs/operations/call-import-sharding-load-test.md):",
            "  • Shard CPU ≤ 75%  • Catalog CPU ≤ 50%  • No pool_timeout / connection exhaustion",
        ],
    )

    _add_content_slide(
        prs,
        "DB Throttling Risks at Scale",
        [
            "Per-pod pool exhaustion: concurrency > catalog pool → SQLAlchemy pool timeout",
            "RDS max_connections: 19 processes × 16/shard × 4 shards = high aggregate without multiplexing",
            "Catalog CPU hotspot: every row reads eval config/metrics from catalog during dispatch + scoring",
            "Hot shard: uneven hash distribution (usually OK for 10k imports with 20 slices across 4 shards)",
            "Write IOPS on shards: concurrent status + metric_scores JSON updates under high inflight",
            "Mitigation now: match concurrency to pool, cap max replicas, monitor dispatch diagnostics",
        ],
    )

    # --- KEDA ---
    _add_section_slide(prs, "8. Kubernetes Autoscaling (KEDA)")

    _add_content_slide(
        prs,
        "KEDA Strategy",
        [
            "minReplicaCount: 8 (never scale below — customer requirement)",
            "maxReplicaCount: 16 (20 only with PgBouncer or larger RDS max_connections)",
            "Do NOT scale on Celery queue depth alone — fair dispatch keeps queues small while Postgres backlog grows",
            "Primary signals (Prometheus):",
            "  • sum(efficientai_eval_rows_pending) — org-wide Postgres backlog",
            "  • eval_inflight_global / eval_global_limit > 0.85 AND pending > 0",
            "  • sum(celery_queue_length{queue=~\"imports|diarization|evaluations\"}) — burst detector",
            "Scale-down: slow (5–10 min stabilization); terminationGracePeriodSeconds: 900 (15 min tasks)",
        ],
    )

    _add_content_slide(
        prs,
        "Observability Stack",
        [
            "Prometheus scrapes: FastAPI /metrics, celery-exporter, redis-exporter, postgres-exporter",
            "celery-exporter (docker-compose.observability.yml): queue lengths, task metrics",
            "Operator endpoint: GET /api/v1/call-imports/dispatch-diagnostics (admin)",
            "  → per-workspace inflight, pending_dispatch_rows, job inflight, RR cursors, at_capacity flags",
            "Grafana dashboards: compare workspace A vs B inflight during dual 10k runs",
            "Alerts: global_at_capacity + high pending; catalog/shard CPU; pool timeout in worker logs",
        ],
    )

    # --- PgBouncer ---
    _add_section_slide(prs, "9. Future Scope: PgBouncer")

    _add_content_slide(
        prs,
        "Why PgBouncer (Next Phase)",
        [
            "Current constraint: each K8s pod opens pool_size + max_overflow connections per catalog + each shard",
            "Without multiplexing: scaling to 20 pods × 28 threads requires 500+ catalog and 300+ per shard connections",
            "PgBouncer sits between app pods and RDS — many client connections, fewer server connections",
            "Enables: higher concurrency per pod, more replicas, higher inflight limits (560 global / 280 per workspace)",
            "Target outcome: dual 10k in ~30 min at max scale without pool_timeout throttling",
        ],
    )

    _add_content_slide(
        prs,
        "PgBouncer Topology (Recommended)",
        [
            "Deploy one PgBouncer pool per logical database:",
            "  • pgbouncer-catalog → efficientai_catalog RDS",
            "  • pgbouncer-shard-01 … pgbouncer-shard-04 → each data shard RDS",
            "Alternative: single PgBouncer instance with multiple database entries (simpler ops, shared process)",
            "App DATABASE_URL / shard URLs point to PgBouncer service DNS, not RDS directly",
            "K8s: PgBouncer as StatefulSet or Helm chart; sidecar pattern generally NOT recommended here",
        ],
    )

    _add_table_slide(
        prs,
        "PgBouncer Settings (Starting Point)",
        ["Parameter", "Catalog pool", "Each data shard pool"],
        [
            ["pool_mode", "transaction", "transaction"],
            ["default_pool_size", "80–120", "50–80"],
            ["max_client_conn", "2000", "2000"],
            ["server_idle_timeout", "600", "600"],
            ["App pool_size (per pod)", "5–8", "5–8"],
            ["App max_overflow", "8–12", "8–12"],
            ["App concurrency (with PgBouncer)", "24–28", "24–28"],
        ],
    )

    _add_content_slide(
        prs,
        "PgBouncer + Scaling Profile (Future)",
        [
            "After PgBouncer rollout, target enterprise profile:",
            "  • 16–20 worker-imports replicas × 28 concurrency = 448–560 thread capacity",
            "  • eval_global = 560, eval_workspace = 280 (÷ 2 workspaces), eval_job = 280",
            "  • eval_fair_dispatch_batch_size = 280",
            "  • import_global = 128–160, import_workspace = 64–80",
            "Expected: 20k rows (dual 10k) in ~30 min @ ~24 s/row with full pod scale",
            "Prerequisite: load-test scenario D (25k rows) pass — shard CPU ≤ 75%, catalog ≤ 50%",
        ],
    )

    _add_content_slide(
        prs,
        "PgBouncer Caveats for SQLAlchemy",
        [
            "Use transaction pooling — compatible with short ORM transactions in workers",
            "Avoid session-level features across transactions: TEMP tables, advisory locks, SET per session",
            "SQLAlchemy: pool_pre_ping=True (already enabled), keep app pools small — let PgBouncer multiplex",
            "Do NOT set app pool_size equal to concurrency without PgBouncer; with PgBouncer, small app pools are correct",
            "Migrate staging first: run call-import-sharding-load-test scenarios A–D before production cutover",
            "Rollback plan: keep direct RDS URLs in config; switch DNS/URL to revert",
        ],
    )

    # --- Next steps ---
    _add_section_slide(prs, "10. Next Steps")

    _add_content_slide(
        prs,
        "Implementation Roadmap",
        [
            "Phase 1 (now, no PgBouncer): Apply recommended config; set eval_job=128; concurrency=16; keep import 96/48",
            "Phase 2: Expose pending-row + inflight Prometheus metrics; tune KEDA on org-wide backlog",
            "Phase 3: Run load-test scenarios A–D on staging; validate dispatch diagnostics during dual 10k",
            "Phase 4: Deploy PgBouncer per catalog + shard; retune to 560/280 enterprise inflight profile",
            "Phase 5: Re-run dual 10k sign-off; document RDS instance sizing and connection budgets",
        ],
    )

    _add_content_slide(
        prs,
        "Key Takeaways",
        [
            "Postgres holds durable row state; Redis enforces fair concurrency — they work together, not interchangeably",
            "Sharding spreads row I/O; catalog remains a shared metadata hub — size and pool it carefully",
            "eval_job_inflight_limit was the hidden bottleneck — not global limit or pod count",
            "Dual 10k across workspaces: set workspace limits = global ÷ 2 for both eval AND import",
            "Queue-length autoscaling alone fails — scale on Postgres pending rows + inflight saturation",
            "PgBouncer unlocks the next tier (560 inflight, 20 pods, dual 10k in ~30 min) safely",
        ],
    )

    _add_title_slide(
        prs,
        "Questions?",
        "EfficientAI Call Import Architecture  |  docs/operations/call-import-sharding-load-test.md",
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build_presentation()
    print(f"Wrote {path}")
